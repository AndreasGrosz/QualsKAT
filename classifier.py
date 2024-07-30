import os
import nltk
import torch
import numpy as np
import pandas as pd
from datasets import Dataset, DatasetDict
from huggingface_hub import HfApi
from requests.exceptions import HTTPError
from transformers import AutoModelForSequenceClassification, TrainingArguments, Trainer, AutoTokenizer
from transformers import DataCollatorWithPadding
import evaluate
from sklearn.preprocessing import LabelEncoder
from nltk.tokenize import sent_tokenize, word_tokenize
import argparse
import docx
from pypdf import PdfReader
import pptx
import xlrd
import logging
import time
from bs4 import BeautifulSoup
import striprtf.striprtf as striprtf
from tqdm import tqdm

# Konfiguriere Logging
logging.basicConfig(filename='file_errors.log', level=logging.ERROR,
                    format='%(asctime)s %(levelname)s:%(message)s')

# Set NLTK data path
nltk.data.path.append("./nltk_data")

def check_hf_token():
    hf_token = os.getenv("HUGGINGFACE_HUB_TOKEN")
    if hf_token is None:
        raise ValueError("Hugging Face API token is missing. Please set the environment variable 'HUGGINGFACE_HUB_TOKEN'.")

    api = HfApi()
    try:
        api.whoami(token=hf_token)
        print("Hugging Face API token is valid.")
    except HTTPError as e:
        raise ValueError("Invalid Hugging Face API token. Please check the token and try again.")

def extract_text_from_file(file_path):
    encodings = ['utf-8', 'latin-1', 'iso-8859-1', 'windows-1252', 'iso-8859-15', 'mac_roman']

    if file_path.endswith('.txt'):
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError as e:
                logging.error(f"UnicodeDecodeError for file {file_path} with encoding {encoding}: {e}")
        logging.error(f"Unable to read the .txt file {file_path} with supported encodings.")
        return None
    elif file_path.endswith('.docx'):
        try:
            doc = docx.Document(file_path)
            return ' '.join([para.text for para in doc.paragraphs])
        except Exception as e:
            logging.error(f"Error reading .docx file {file_path}: {e}")
            return None
    elif file_path.endswith('.pdf'):
        try:
            reader = PdfReader(file_path)
            return ' '.join([page.extract_text() for page in reader.pages])
        except Exception as e:
            logging.error(f"Error reading .pdf file {file_path}: {e}")
            return None
    elif file_path.endswith('.pptx'):
        try:
            ppt = pptx.Presentation(file_path)
            return ' '.join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, 'text')])
        except Exception as e:
            logging.error(f"Error reading .pptx file {file_path}: {e}")
            return None
    elif file_path.endswith('.xlsx') or file_path.endswith('.xls'):
        try:
            workbook = xlrd.open_workbook(file_path)
            return ' '.join([sheet.cell_value(row, col)
                             for sheet in workbook.sheets()
                             for row in range(sheet.nrows)
                             for col in range(sheet.ncols)
                             if sheet.cell_value(row, col)])
        except Exception as e:
            logging.error(f"Error reading Excel file {file_path}: {e}")
            return None
    elif file_path.endswith('.rtf'):
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                rtf_text = file.read()
                return striprtf.striprtf(rtf_text)
        except Exception as e:
            logging.error(f"Error reading .rtf file {file_path}: {e}")
            return None
    elif file_path.endswith('.html') or file_path.endswith('.htm'):
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file, 'html.parser')
                return soup.get_text()
        except Exception as e:
            logging.error(f"Error reading HTML file {file_path}: {e}")
            return None
    else:
        logging.error(f"Unsupported file type: {file_path}")
        return None

def get_files_and_categories(root_dir, test_mode=False):
    files_and_categories = []
    supported_formats = ['.txt', '.doc', '.docx', '.rtf', '.pdf', '.html', '.htm']

    for root, dirs, files in os.walk(root_dir):
        for file in files:
            if any(file.endswith(fmt) for fmt in supported_formats):
                file_path = os.path.join(root, file)
                categories = os.path.relpath(root, root_dir).split(os.sep)

                if categories[0] == "LRH":
                    # LRH mit Unterkategorien
                    category = f"LRH-{'-'.join(categories[1:])}"
                elif categories[0] == "unknown":
                    category = "unknown"
                else:
                    # Potenzielle Ghostwriter
                    category = f"Ghostwriter-{categories[0]}"

                files_and_categories.append((file_path, category))

                if test_mode and len(files_and_categories) >= 10:
                    return files_and_categories
    return files_and_categories

def collect_category_sizes(root_dir):
    category_sizes = {}
    files_and_categories = get_files_and_categories(root_dir)

    for file_path, category in tqdm(files_and_categories, desc="Sammle Kategoriegrößen"):
        text = extract_text_from_file(file_path)
        if text:
            size = len(text.encode('utf-8'))  # Größe in Bytes
            if category in category_sizes:
                category_sizes[category] += size
            else:
                category_sizes[category] = size

    return category_sizes

def print_category_sizes(category_sizes):
    print("Kategorien und ihre Größen:")

    lrh_total = 0
    ghostwriter_total = 0
    unknown_total = 0

    for category, size in sorted(category_sizes.items()):
        print(f"{category}: {size} Bytes")
        if category.startswith("LRH"):
            lrh_total += size
        elif category.startswith("Ghostwriter"):
            ghostwriter_total += size
        elif category == "unknown":
            unknown_total += size

    print("\nZusammenfassung:")
    print(f"LRH Gesamt: {lrh_total} Bytes")
    print(f"Ghostwriter Gesamt: {ghostwriter_total} Bytes")
    print(f"Unbekannt Gesamt: {unknown_total} Bytes")
    print(f"Gesamtzahl der Kategorien: {len(category_sizes)}")
    print(f"Gesamtgröße aller Kategorien: {sum(category_sizes.values())} Bytes")
    print()

def create_dataset(root_dir, test_mode=False):
    files_and_categories = get_files_and_categories(root_dir, test_mode)
    texts = []
    all_categories = []

    for file_path, category in files_and_categories:
        text = extract_text_from_file(file_path)
        if text:
            if test_mode:
                text = text[:10000]  # Limit to first 10000 characters in test mode
            texts.append(text)
            all_categories.append(category)

    le = LabelEncoder()
    numeric_categories = le.fit_transform(all_categories)

    dataset_dict = {
        'text': texts,
        'labels': numeric_categories
    }

    return Dataset.from_dict(dataset_dict), le

def preprocess_function(examples):
    tokenizer = AutoTokenizer.from_pretrained("distilbert-base-uncased")
    return tokenizer(examples["text"], truncation=True)

accuracy = evaluate.load("accuracy")

def compute_metrics(eval_pred):
    predictions, labels = eval_pred
    predictions = np.argmax(predictions, axis=1)
    return accuracy.compute(predictions=predictions, references=labels)

def predict_top_n(trainer, tokenizer, text, le, n=3):
    inputs = tokenizer(text, return_tensors="pt", truncation=True, padding=True)
    with torch.no_grad():
        logits = trainer.model(**inputs).logits
    probabilities = torch.nn.functional.softmax(logits, dim=-1)
    top_n_prob, top_n_indices = torch.topk(probabilities, n)
    results = []
    for prob, idx in zip(top_n_prob[0], top_n_indices[0]):
        category = le.inverse_transform([idx.item()])[0]
        results.append((category, prob.item()))
    return results

def main(args):
    start_time = time.time()
    print(f"Script started at: {time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(start_time))}")

    check_hf_token()

    root_dir = 'documents'
    test_mode = args.test_mode

    # Sammeln und Ausgeben der Kategoriegrößen
    category_sizes = collect_category_sizes(root_dir)
    print_category_sizes(category_sizes)

    dataset, le = create_dataset(root_dir, test_mode)

    dataset = dataset.train_test_split(test_size=0.2)
    test_valid = dataset['test'].train_test_split(test_size=0.5)
    dataset_dict = DatasetDict({
        'train': dataset['train'],
        'test': test_valid['test'],
        'validation': test_valid['train']
    })

    tokenizer = AutoTokenizer.from_pretrained("distilbert-base-uncased")
    tokenized_dataset = dataset_dict.map(lambda examples: tokenizer(examples['text'], truncation=True, padding='max_length'), batched=True)

    num_labels = len(le.classes_)
    model = AutoModelForSequenceClassification.from_pretrained(
        "distilbert-base-uncased",
        num_labels=num_labels
    )

    training_args = TrainingArguments(
        output_dir='Multi_Class_Classifier',
        learning_rate=2e-5,
        per_device_train_batch_size=16,
        per_device_eval_batch_size=16,
        num_train_epochs=3,
        weight_decay=0.01,
        eval_strategy="epoch",
        save_strategy="epoch",
        load_best_model_at_end=True,
    )

    trainer = Trainer(
        model=model,
        args=training_args,
        train_dataset=tokenized_dataset['train'],
        eval_dataset=tokenized_dataset['validation'],
        tokenizer=tokenizer,
        compute_metrics=compute_metrics
    )

    trainer.train()

    results = trainer.evaluate(tokenized_dataset['test'])
    print("Test results:", results)

    example_texts = [
        "The tech of auditing involves the use of an E-meter to locate areas of spiritual distress.",
        "In the 1950s, LRH gave lectures on Dianetics and the fundamentals of Scientology.",
        "David Mayo was known for his contributions to the upper levels of Scientology tech.",
        "This is an unknown text that doesn't clearly belong to any specific author.",
    ]

    for text in example_texts:
        top_predictions = predict_top_n(trainer, tokenizer, text, le)
        print(f"Text: '{text[:50]}...'")
        for category, confidence in top_predictions:
            print(f"  {category}: {confidence:.2f}")
        print()

    end_time = time.time()
    print(f"Script ended at: {time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(end_time))}")
    print(f"Total runtime: {end_time - start_time:.2f} seconds")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Multi-Format Document Classification')
    parser.add_argument('-approach', choices=['discriminative'], required=True)
    parser.add_argument('-test_mode', action='store_true', help='Run in test mode with limited data')
    args = parser.parse_args()
    main(args)
